"""Build executive PowerPoint deck from salary sacrifice analysis."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path
import copy

OUT = Path("outputs/2026-04-29_salary_sacrifice_report.pptx")
CHARTS = Path("outputs")

# Brand colours
NAVY   = RGBColor(0x1A, 0x2E, 0x4A)
TEAL   = RGBColor(0x00, 0x87, 0x8A)
GOLD   = RGBColor(0xF0, 0xA8, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF4, 0xF5, 0xF7)
DGREY  = RGBColor(0x44, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def fill_bg(slide, colour: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_rect(slide, l, t, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, colour=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.color.rgb = colour
    run.font.name    = "Calibri"
    return txb


def add_bullet_box(slide, bullets, l, t, w, h,
                   font_size=16, colour=DGREY, spacing_after=6):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = f"  {bullet}"
        run.font.size      = Pt(font_size)
        run.font.color.rgb = colour
        run.font.name      = "Calibri"
    return txb


# ── Slide helpers ─────────────────────────────────────────────────────────────

def slide_title_page(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, NAVY)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, TEAL)

    # Gold rule
    add_rect(slide, Inches(0.45), Inches(3.55), Inches(8), Inches(0.06), GOLD)

    add_text(slide, "Salary Sacrifice Behaviour",
             Inches(0.45), Inches(1.6), Inches(9), Inches(1.1),
             font_size=40, bold=True, colour=WHITE)
    add_text(slide, "Across Age Bands",
             Inches(0.45), Inches(2.6), Inches(9), Inches(0.9),
             font_size=40, bold=True, colour=WHITE)

    add_text(slide, "Superannuation Member Analysis  |  Snapshot: 31 March 2026",
             Inches(0.45), Inches(3.75), Inches(10), Inches(0.5),
             font_size=14, colour=RGBColor(0xB0, 0xC4, 0xDE))
    add_text(slide, "Prepared: 29 April 2026  |  Analytics & Insights",
             Inches(0.45), Inches(4.2), Inches(10), Inches(0.5),
             font_size=14, colour=RGBColor(0xB0, 0xC4, 0xDE))

    # Bottom strip
    add_rect(slide, Inches(0), Inches(6.9), SLIDE_W, Inches(0.6), TEAL)
    add_text(slide, "CONFIDENTIAL — FOR INTERNAL USE ONLY",
             Inches(0.3), Inches(6.92), Inches(10), Inches(0.4),
             font_size=11, colour=WHITE, align=PP_ALIGN.LEFT)


def slide_exec_summary(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, LGREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), NAVY)
    add_text(slide, "Executive Summary",
             Inches(0.4), Inches(0.25), Inches(12), Inches(0.7),
             font_size=26, bold=True, colour=WHITE)

    bullets = [
        "Participation rises sharply with age: 9% of 18–24 year-olds salary sacrifice vs 64% of the 55–64 cohort.",
        "Age drives participation independently of salary — it is not simply that younger members earn less.",
        "Once a member participates, the dollar amount is nearly identical across all age bands (~$6,000–$7,400 median).",
        "Younger participants sacrifice a higher share of their salary (~7.0%) than older ones (~5.1%).",
        "Biggest opportunity: converting 718 non-participating 35–44 year-olds with meaningful earning power.",
    ]

    y = Inches(1.3)
    for i, b in enumerate(bullets):
        # Number circle
        circ = slide.shapes.add_shape(9, Inches(0.35), y + Inches(0.05),
                                       Inches(0.42), Inches(0.42))
        circ.fill.solid(); circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        tf = circ.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1); run.font.size = Pt(13)
        run.font.bold = True; run.font.color.rgb = WHITE
        run.font.name = "Calibri"

        add_text(slide, b, Inches(0.9), y, Inches(11.9), Inches(0.55),
                 font_size=15, colour=DGREY)
        y += Inches(0.98)

    # Source note
    add_rect(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), NAVY)
    add_text(slide, "Source: super_members_2026Q1.csv  |  n=5,000 members (25 duplicates removed)",
             Inches(0.3), Inches(7.07), Inches(12), Inches(0.38),
             font_size=10, colour=RGBColor(0xB0, 0xC4, 0xDE))


def slide_with_chart(prs, title, subtitle, chart_file, annotation, note=""):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
             font_size=22, bold=True, colour=WHITE)
    add_text(slide, subtitle,
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.4),
             font_size=13, colour=RGBColor(0xB0, 0xC4, 0xDE))

    # Chart image — left column
    chart_path = CHARTS / chart_file
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path),
                                  Inches(0.3), Inches(1.2),
                                  Inches(8.5), Inches(5.6))

    # Annotation panel — right column
    add_rect(slide, Inches(9.0), Inches(1.2), Inches(4.0), Inches(5.6), LGREY)
    add_rect(slide, Inches(9.0), Inches(1.2), Inches(0.07), Inches(5.6), TEAL)
    add_text(slide, annotation,
             Inches(9.2), Inches(1.35), Inches(3.7), Inches(5.2),
             font_size=13, colour=DGREY, wrap=True)

    if note:
        add_rect(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), LGREY)
        add_text(slide, note, Inches(0.3), Inches(7.07), Inches(12.7), Inches(0.38),
                 font_size=10, colour=RGBColor(0x88, 0x88, 0x88))


def slide_two_col(prs, title, subtitle, left_title, left_bullets,
                  right_title, right_bullets, accent_right=TEAL):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
             font_size=22, bold=True, colour=WHITE)
    add_text(slide, subtitle,
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.4),
             font_size=13, colour=RGBColor(0xB0, 0xC4, 0xDE))

    # Left panel
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(6.1), Inches(5.7), LGREY)
    add_rect(slide, Inches(0.3), Inches(1.2), Inches(0.07), Inches(5.7), TEAL)
    add_text(slide, left_title,
             Inches(0.55), Inches(1.3), Inches(5.7), Inches(0.45),
             font_size=14, bold=True, colour=NAVY)
    add_bullet_box(slide, left_bullets,
                   Inches(0.55), Inches(1.85), Inches(5.7), Inches(4.8),
                   font_size=13, colour=DGREY)

    # Right panel
    add_rect(slide, Inches(6.9), Inches(1.2), Inches(6.1), Inches(5.7), LGREY)
    add_rect(slide, Inches(6.9), Inches(1.2), Inches(0.07), Inches(5.7), accent_right)
    add_text(slide, right_title,
             Inches(7.1), Inches(1.3), Inches(5.7), Inches(0.45),
             font_size=14, bold=True, colour=NAVY)
    add_bullet_box(slide, right_bullets,
                   Inches(7.1), Inches(1.85), Inches(5.7), Inches(4.8),
                   font_size=13, colour=DGREY)


def slide_opportunities(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, WHITE)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text(slide, "Implications & Opportunities",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
             font_size=22, bold=True, colour=WHITE)
    add_text(slide, "Priority actions to increase salary sacrifice participation",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.4),
             font_size=13, colour=RGBColor(0xB0, 0xC4, 0xDE))

    cards = [
        ("18–34\n1,125 non-participants",
         TEAL,
         "Salary milestone triggers.\nTarget members crossing ~$85K salary threshold with proactive comms. For very young low-income members, focus on habit formation."),
        ("35–44\n718 non-participants",
         GOLD,
         "Biggest absolute opportunity.\nPeer comparison messaging + default opt-in nudges. Half their cohort already participates — social proof is actionable."),
        ("55–64\n273 non-participants",
         RGBColor(0x6A, 0x0D, 0xAD),
         "Investigate before assuming disengagement.\nMany have high salaries. Check for alternative strategies (personal deductible contributions). Member service call may outperform digital nudge."),
        ("65+ Participants\nn=175",
         NAVY,
         "Concessional cap guidance.\nMany are near the $30K cap. Proactive guidance on carry-forward unused concessional contributions could unlock additional savings."),
    ]

    x = Inches(0.3)
    for label, colour, body in cards:
        add_rect(slide, x, Inches(1.25), Inches(3.0), Inches(5.9), LGREY)
        add_rect(slide, x, Inches(1.25), Inches(3.0), Inches(0.7), colour)
        add_text(slide, label,
                 x + Inches(0.1), Inches(1.3), Inches(2.8), Inches(0.65),
                 font_size=12, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)
        add_text(slide, body,
                 x + Inches(0.12), Inches(2.05), Inches(2.78), Inches(4.9),
                 font_size=12, colour=DGREY, wrap=True)
        x += Inches(3.25)


def slide_caveats(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    fill_bg(slide, LGREY)

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text(slide, "Caveats & Limitations",
             Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
             font_size=22, bold=True, colour=WHITE)
    add_text(slide, "Considerations for interpreting these findings",
             Inches(0.4), Inches(0.65), Inches(12), Inches(0.4),
             font_size=13, colour=RGBColor(0xB0, 0xC4, 0xDE))

    caveats = [
        "Cross-sectional snapshot only (31/03/2026) — causality cannot be established; members not sacrificing may have done so previously.",
        "Salary nulls: 100 members (2%) missing salary data, excluded from salary analyses. May bias salary-stratified results if non-random.",
        "18–24 participant group is small (n=31) — treat conditional amounts for this band as indicative only.",
        "No adviser relationship data — members with a financial adviser are likely more inclined to salary sacrifice; an omitted variable.",
        "Concessional cap not modelled — the $30K annual cap (inclusive of employer SG) was not applied; explains plateaus in older high-earning bands.",
        "Regression R² = 9.5% — age, salary, and tenure explain less than 10% of participation variation; significant unexplained factors remain.",
    ]

    y = Inches(1.3)
    for c in caveats:
        add_rect(slide, Inches(0.35), y + Inches(0.1), Inches(0.06), Inches(0.32), GOLD)
        add_text(slide, c,
                 Inches(0.6), y, Inches(12.3), Inches(0.6),
                 font_size=13, colour=DGREY)
        y += Inches(0.88)

    add_rect(slide, Inches(0), Inches(7.05), SLIDE_W, Inches(0.45), NAVY)
    add_text(slide, "Data source: super_members_2026Q1.csv  |  n=5,000 members  |  Analysis: Analytics & Insights, 29 April 2026",
             Inches(0.3), Inches(7.07), Inches(12.7), Inches(0.38),
             font_size=10, colour=RGBColor(0xB0, 0xC4, 0xDE))


# ── Build deck ────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()

    slide_title_page(prs)

    slide_exec_summary(prs)

    slide_with_chart(
        prs,
        title="The Age Gap Is a Participation Problem",
        subtitle="Participation rate climbs from 9% (18–24) to 64% (55–64) — conditional amounts are flat",
        chart_file="2026-04-29_participation_rate_by_age_band.png",
        annotation=(
            "KEY FINDING\n\n"
            "9.2% of 18–24 year-olds make any salary sacrifice.\n\n"
            "63.5% of 55–64 year-olds do.\n\n"
            "The gap is who opts in — not how much they contribute once they do.\n\n"
            "Overall fund participation: 48.7%\n(2,434 of 5,000 members)"
        ),
        note="Source: super_members_2026Q1.csv  |  n=5,000 members"
    )

    slide_with_chart(
        prs,
        title="Conditional Amounts Are Consistent Across Age Bands",
        subtitle="Among participants, median annual salary sacrifice ranges narrowly from $6,024 to $7,396",
        chart_file="2026-04-29_conditional_amount_by_age_band.png",
        annotation=(
            "WHAT THIS MEANS\n\n"
            "Once a member opts in, the amount they contribute does not vary significantly by age.\n\n"
            "Median: ~$6,000–$7,400 across all bands.\n\n"
            "The 18–24 median ($7,396) looks high but is based on only 31 participants — treat as indicative.\n\n"
            "Strategy implication: focus on activation, not on lifting contribution amounts."
        ),
        note="Participants only (n=2,434)  |  Source: super_members_2026Q1.csv"
    )

    slide_with_chart(
        prs,
        title="Salary Matters — But Age Has the Stronger Effect",
        subtitle="Within each salary quartile, older members still participate at higher rates",
        chart_file="2026-04-29_participation_by_age_salary_quartile.png",
        annotation=(
            "REGRESSION RESULTS\n\n"
            "Salary alone explains 3.8% of participation variation (R²).\n\n"
            "Adding age rank + tenure raises this to 9.5%.\n\n"
            "Standardised coefficients:\n"
            "  Age rank:   0.110\n"
            "  Log salary: 0.048\n"
            "  Tenure:     0.028\n\n"
            "Age is ~2× more predictive than salary.\n\n"
            "Note: In the 55–64 band, Q4 earners participate less than Q2/Q3 — likely concessional cap constraints."
        ),
        note="OLS linear probability model  |  Source: super_members_2026Q1.csv"
    )

    slide_with_chart(
        prs,
        title="Younger Participants Sacrifice a Higher Share of Their Salary",
        subtitle="Sacrifice as % of salary declines with age — likely due to the concessional cap becoming binding",
        chart_file="2026-04-29_sacrifice_pct_salary_by_age.png",
        annotation=(
            "% OF SALARY SACRIFICED\n(participants only)\n\n"
            "18–24:  7.0%\n"
            "25–34:  6.6%\n"
            "35–44:  6.3%\n"
            "45–54:  6.1%\n"
            "55–64:  5.7%\n"
            "65+:    5.1%\n\n"
            "As salaries rise with age, the $30K concessional cap (including employer SG) limits the effective sacrifice rate.\n\n"
            "Younger participants are stretching proportionally further on lower incomes."
        ),
        note="Participants only (n=2,434)  |  Concessional cap not modelled  |  Source: super_members_2026Q1.csv"
    )

    slide_two_col(
        prs,
        title="Other Factors: What Matters and What Doesn't",
        subtitle="Tenure and investment option show modest signals; digital engagement does not predict participation",
        left_title="Factors with some signal",
        left_bullets=[
            "Tenure: participants have slightly higher tenure in every band (e.g. 5.2 vs 4.3 yrs in 18–24). Effect is small but consistent (β=0.028).",
            "Investment option: High Growth members in 18–24 band participate at 16.2% vs 7.1% for Balanced — but this likely reflects engaged-member selection bias.",
            "Gender: small and inconsistent gap (males 11.7% vs females 7.1% in 18–24). Effect reverses partially in older bands.",
        ],
        right_title="Factors that don't predict participation",
        right_bullets=[
            "App installed: members with the app installed participate at lower rates in several bands (e.g. 25.9% vs 30.3% in 25–34). No positive signal.",
            "Login recency: median days since last login nearly identical for participants and non-participants across all bands (typically 29–39 days for both).",
            "Digital engagement, as measured in this dataset, does not explain salary sacrifice behaviour.",
        ],
        accent_right=GOLD,
    )

    slide_opportunities(prs)

    slide_caveats(prs)

    prs.save(str(OUT))
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    build()
