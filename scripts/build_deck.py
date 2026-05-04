"""
Parametric slide deck builder — consultant-quality output for any analysis.

Usage:
    python scripts/build_deck.py outputs/findings_config.json

Config format: see DeckConfig / SlideContent at bottom of file.
"""

from __future__ import annotations
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
import json
import sys
import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design system ────────────────────────────────────────────────────────────

class C:
    """Centralised colour palette — every colour used in the deck lives here."""
    NAVY       = RGBColor(0x1B, 0x35, 0x57)   # Primary — headlines, backgrounds
    TEAL       = RGBColor(0x0A, 0x8A, 0x8F)   # Accent 1 — rules, highlights
    AMBER      = RGBColor(0xE8, 0xA0, 0x20)   # Accent 2 — alerts, opportunity
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    CHARCOAL   = RGBColor(0x2C, 0x3E, 0x50)   # Body text
    LGREY      = RGBColor(0xF7, 0xF8, 0xFA)   # Slide backgrounds
    MGREY      = RGBColor(0xE5, 0xE8, 0xEC)   # Dividers
    SLATE      = RGBColor(0x8A, 0x99, 0xAA)   # Metadata, source notes
    POSITIVE   = RGBColor(0x18, 0x7A, 0x4E)   # Green — above benchmark
    NEGATIVE   = RGBColor(0xC0, 0x39, 0x2B)   # Red — below benchmark / risk
    COVER_SUB  = RGBColor(0x8A, 0xA8, 0xC8)   # Cover subtitle text

FONT = "Calibri"

W = Inches(13.33)
H = Inches(7.5)

# ── Low-level primitives ─────────────────────────────────────────────────────

def _blank_layout(prs: Presentation):
    return prs.slide_layouts[6]


def _bg(slide, colour: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _rect(slide, l, t, w, h, fill: RGBColor, *, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if not line:
        s.line.fill.background()
    return s


def _txt(slide, text: str, l, t, w, h, *,
         size: float = 14, bold: bool = False,
         colour: RGBColor = C.CHARCOAL,
         align: PP_ALIGN = PP_ALIGN.LEFT,
         wrap: bool = True,
         italic: bool = False) -> None:
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = colour
    run.font.name = FONT


def _bullets(slide, items: list[str], l, t, w, h, *,
             size: float = 14, colour: RGBColor = C.CHARCOAL,
             spacing_after: float = 7, bold_prefix: bool = False):
    """Render a list of strings as bullets with an em-dash character."""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        # Split bold prefix (text before first '|') if present
        if bold_prefix and '|' in item:
            prefix, rest = item.split('|', 1)
            run.text = f"  {prefix.strip()}"
            run.font.bold = True
            run.font.color.rgb = C.NAVY
            run.font.name = FONT
            run.font.size = Pt(size)
            run2 = p.add_run()
            run2.text = f"  {rest.strip()}"
            run2.font.bold = False
            run2.font.color.rgb = colour
            run2.font.name = FONT
            run2.font.size = Pt(size)
        else:
            run.text = f"  –  {item}"
            run.font.size = Pt(size)
            run.font.color.rgb = colour
            run.font.name = FONT


def _footer(slide, source_note: str):
    """Standard footer with a hairline rule and source text."""
    _rect(slide, Inches(0), Inches(7.2), W, Inches(0.02), C.MGREY)
    _txt(slide, source_note,
         Inches(0.4), Inches(7.25), Inches(12.5), Inches(0.25),
         size=8.5, colour=C.SLATE, italic=True)


def _slide_number(slide, n: int, total: int):
    _txt(slide, f"{n} / {total}",
         Inches(12.6), Inches(7.25), Inches(0.7), Inches(0.25),
         size=8.5, colour=C.SLATE, align=PP_ALIGN.RIGHT)


def _header_band(slide, title: str, subtitle: str = ""):
    """Standard dark header band used on content slides."""
    _rect(slide, Inches(0), Inches(0), W, Inches(1.0), C.NAVY)
    _txt(slide, title,
         Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.58),
         size=22, bold=True, colour=C.WHITE)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.4), Inches(0.64), Inches(12.5), Inches(0.34),
             size=12, colour=C.COVER_SUB)


def _existing_asset(path: str | Path, *, label: str) -> Path:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(f"{label} not found: {p}")
    return p


# ── Slide constructors ───────────────────────────────────────────────────────

def slide_cover(prs: Presentation, cfg: DeckConfig) -> None:
    """Impact cover — dark background, large title, minimal."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.NAVY)

    # Left accent stripe
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, C.TEAL)

    # Teal rule below title
    _rect(slide, Inches(0.5), Inches(3.7), Inches(7.5), Inches(0.055), C.TEAL)

    # Title — split across two lines if needed
    lines = cfg.title.split('\n') if '\n' in cfg.title else [cfg.title]
    y = Inches(1.8)
    for line in lines:
        _txt(slide, line, Inches(0.5), y, Inches(10.5), Inches(1.1),
             size=42, bold=True, colour=C.WHITE)
        y += Inches(0.95)

    # Metadata below rule
    _txt(slide, cfg.subtitle,
         Inches(0.5), Inches(3.9), Inches(10.5), Inches(0.5),
         size=14, colour=C.COVER_SUB)

    meta_parts = []
    if cfg.organisation:
        meta_parts.append(cfg.organisation)
    if cfg.prepared_date:
        meta_parts.append(cfg.prepared_date)
    if cfg.prepared_by:
        meta_parts.append(f"Prepared by: {cfg.prepared_by}")
    _txt(slide, "  ·  ".join(meta_parts),
         Inches(0.5), Inches(4.45), Inches(10.5), Inches(0.4),
         size=12, colour=C.COVER_SUB)

    # Bottom confidential bar
    _rect(slide, Inches(0), Inches(7.05), W, Inches(0.45), C.TEAL)
    _txt(slide, cfg.classification,
         Inches(0.45), Inches(7.07), Inches(12), Inches(0.38),
         size=10, colour=C.WHITE)


def slide_exec_summary(prs: Presentation, bullets: list[str],
                       source_note: str, slide_num: int, total: int) -> None:
    """Executive summary — numbered callouts with large visual weight."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.LGREY)
    _header_band(slide, "Executive Summary",
                 "Key findings — read this slide first")

    accent_colours = [C.TEAL, C.TEAL, C.AMBER, C.TEAL, C.TEAL]
    y = Inches(1.12)
    for i, bullet in enumerate(bullets[:5]):
        accent = accent_colours[i % len(accent_colours)]

        # Number pill
        pill = slide.shapes.add_shape(
            9,  # Rounded rectangle
            Inches(0.35), y + Inches(0.04),
            Inches(0.42), Inches(0.42)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = accent
        pill.line.fill.background()
        tf = pill.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(i + 1)
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = C.WHITE
        run.font.name = FONT

        # Left accent rule
        _rect(slide, Inches(0.92), y + Inches(0.07), Inches(0.035), Inches(0.3), accent)

        # Bullet text
        _txt(slide, bullet,
             Inches(1.05), y, Inches(12.0), Inches(0.53),
             size=14.5, colour=C.CHARCOAL)

        y += Inches(1.05)

    _footer(slide, source_note)
    _slide_number(slide, slide_num, total)


def slide_key_number(prs: Presentation, title: str, subtitle: str,
                     hero_number: str, hero_label: str,
                     supporting_points: list[str],
                     source_note: str, slide_num: int, total: int) -> None:
    """Hero-number slide — one massive statistic, supporting context alongside."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    # Left panel — hero number
    _txt(slide, hero_number,
         Inches(0.4), Inches(1.2), Inches(5.8), Inches(2.8),
         size=96, bold=True, colour=C.NAVY, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(0.4), Inches(3.9), Inches(5.8), Inches(0.04), C.TEAL)
    _txt(slide, hero_label,
         Inches(0.4), Inches(4.0), Inches(5.8), Inches(0.7),
         size=14, colour=C.SLATE, align=PP_ALIGN.CENTER)

    # Right panel — context
    _rect(slide, Inches(6.6), Inches(1.2), Inches(6.4), Inches(5.9), C.LGREY)
    _rect(slide, Inches(6.6), Inches(1.2), Inches(0.055), Inches(5.9), C.TEAL)
    _bullets(slide, supporting_points,
             Inches(6.8), Inches(1.35), Inches(6.0), Inches(5.5),
             size=13.5, colour=C.CHARCOAL, spacing_after=10)

    _footer(slide, source_note)
    _slide_number(slide, slide_num, total)


def slide_chart(prs: Presentation, title: str, subtitle: str,
                chart_path: str, annotation_lines: list[str],
                source_note: str, slide_num: int, total: int) -> None:
    """Chart + annotation — 70/30 split, clean."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    p = _existing_asset(chart_path, label="Chart asset")
    slide.shapes.add_picture(str(p), Inches(0.3), Inches(1.12), Inches(8.8), Inches(5.9))

    # Right annotation panel
    _rect(slide, Inches(9.3), Inches(1.12), Inches(3.7), Inches(5.9), C.LGREY)
    _rect(slide, Inches(9.3), Inches(1.12), Inches(0.055), Inches(5.9), C.AMBER)

    _txt(slide, "KEY INSIGHT",
         Inches(9.5), Inches(1.25), Inches(3.3), Inches(0.3),
         size=9, bold=True, colour=C.AMBER)

    body = "\n\n".join(annotation_lines)
    _txt(slide, body,
         Inches(9.5), Inches(1.6), Inches(3.35), Inches(5.1),
         size=12.5, colour=C.CHARCOAL)

    _footer(slide, source_note)
    _slide_number(slide, slide_num, total)


def slide_two_chart(prs: Presentation, title: str, subtitle: str,
                    left_chart: str, left_label: str,
                    right_chart: str, right_label: str,
                    source_note: str, slide_num: int, total: int) -> None:
    """Two-chart side-by-side layout."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    for path, x, label in [
        (left_chart, Inches(0.25), left_label),
        (right_chart, Inches(6.85), right_label),
    ]:
        p = _existing_asset(path, label="Chart asset")
        slide.shapes.add_picture(str(p), x, Inches(1.12), Inches(6.35), Inches(4.8))
        _rect(slide, x, Inches(6.05), Inches(6.35), Inches(0.03), C.TEAL)
        _txt(slide, label, x, Inches(6.1), Inches(6.35), Inches(0.6),
             size=11, colour=C.SLATE, italic=True, align=PP_ALIGN.CENTER)

    _footer(slide, source_note)
    _slide_number(slide, slide_num, total)


def slide_section_divider(prs: Presentation, section_number: str,
                          section_title: str, description: str = "") -> None:
    """Dark full-bleed section divider."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.NAVY)

    _rect(slide, Inches(0.4), Inches(2.6), Inches(1.2), Inches(0.07), C.TEAL)
    _txt(slide, section_number,
         Inches(0.4), Inches(1.5), Inches(4), Inches(0.8),
         size=48, bold=True, colour=C.TEAL)
    _txt(slide, section_title,
         Inches(0.4), Inches(2.8), Inches(9), Inches(1.2),
         size=34, bold=True, colour=C.WHITE)
    if description:
        _txt(slide, description,
             Inches(0.4), Inches(4.2), Inches(8.5), Inches(0.7),
             size=14, colour=C.COVER_SUB)


def slide_two_col(prs: Presentation, title: str, subtitle: str,
                  left_heading: str, left_items: list[str],
                  right_heading: str, right_items: list[str],
                  left_accent: RGBColor = None, right_accent: RGBColor = None,
                  source_note: str = "", slide_num: int = 0, total: int = 0) -> None:
    """Two-column comparison — What works / What doesn't, Before / After, etc."""
    left_accent = left_accent or C.TEAL
    right_accent = right_accent or C.AMBER
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    for panel_x, accent, heading, items in [
        (Inches(0.3), left_accent, left_heading, left_items),
        (Inches(6.9), right_accent, right_heading, right_items),
    ]:
        _rect(slide, panel_x, Inches(1.12), Inches(6.1), Inches(6.1), C.LGREY)
        _rect(slide, panel_x, Inches(1.12), Inches(0.055), Inches(6.1), accent)
        _txt(slide, heading,
             panel_x + Inches(0.2), Inches(1.2), Inches(5.8), Inches(0.45),
             size=13.5, bold=True, colour=C.NAVY)
        _rect(slide, panel_x + Inches(0.2), Inches(1.73), Inches(5.7), Inches(0.025), C.MGREY)
        _bullets(slide, items,
                 panel_x + Inches(0.2), Inches(1.82), Inches(5.75), Inches(5.1),
                 size=13, colour=C.CHARCOAL, spacing_after=9)

    if source_note:
        _footer(slide, source_note)
    if slide_num:
        _slide_number(slide, slide_num, total)


def slide_action_cards(prs: Presentation, title: str, subtitle: str,
                       cards: list[dict],
                       source_note: str = "", slide_num: int = 0, total: int = 0) -> None:
    """
    Recommendation/opportunity cards — up to 4 cards per slide.
    Each card: {"label": str, "metric": str, "body": str, "accent": RGBColor (optional)}
    """
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    card_w = Inches(3.0)
    gap = Inches(0.17)
    start_x = Inches(0.35)

    accents = [C.TEAL, C.AMBER, C.NAVY, RGBColor(0x6A, 0x0D, 0xAD)]

    for i, card in enumerate(cards[:4]):
        x = start_x + i * (card_w + gap)
        accent = card.get("accent") or accents[i % len(accents)]

        # Card background
        _rect(slide, x, Inches(1.15), card_w, Inches(5.95), C.LGREY)

        # Coloured top bar
        _rect(slide, x, Inches(1.15), card_w, Inches(0.78), accent)

        # Card label and metric in header
        _txt(slide, card["label"],
             x + Inches(0.12), Inches(1.2), card_w - Inches(0.2), Inches(0.38),
             size=11.5, bold=True, colour=C.WHITE)
        if card.get("metric"):
            _txt(slide, card["metric"],
                 x + Inches(0.12), Inches(1.58), card_w - Inches(0.2), Inches(0.32),
                 size=10.5, colour=RGBColor(0xE8, 0xF4, 0xF8), italic=True)

        # Card body
        _txt(slide, card["body"],
             x + Inches(0.14), Inches(2.08), card_w - Inches(0.25), Inches(4.85),
             size=12.5, colour=C.CHARCOAL)

    if source_note:
        _footer(slide, source_note)
    if slide_num:
        _slide_number(slide, slide_num, total)


def slide_table(prs: Presentation, title: str, subtitle: str,
                headers: list[str], rows: list[list[str]],
                col_widths: list[float],
                source_note: str = "", slide_num: int = 0, total: int = 0) -> None:
    """Clean data table slide."""
    from pptx.util import Inches as I_
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.WHITE)
    _header_band(slide, title, subtitle)

    total_w = sum(col_widths)
    tbl_x = Inches(0.4)
    tbl_y = Inches(1.18)
    tbl_w = Inches(total_w)
    tbl_h = Inches(0.45 * (1 + len(rows)))

    table = slide.shapes.add_table(
        1 + len(rows), len(headers), tbl_x, tbl_y, tbl_w, tbl_h
    ).table

    # Header row
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.NAVY
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.bold = True
        run.font.color.rgb = C.WHITE
        run.font.name = FONT
        run.font.size = Pt(11)
        table.columns[c].width = Inches(col_widths[c])

    # Data rows
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = table.cell(r + 1, c)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C.LGREY
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.color.rgb = C.CHARCOAL
            run.font.name = FONT
            run.font.size = Pt(11)

    if source_note:
        _footer(slide, source_note)
    if slide_num:
        _slide_number(slide, slide_num, total)


def slide_caveats(prs: Presentation, caveats: list[str],
                  source_note: str, slide_num: int, total: int) -> None:
    """Caveats and limitations — amber rule accent."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.LGREY)
    _header_band(slide, "Caveats & Limitations",
                 "Factors to consider when interpreting these findings")

    y = Inches(1.18)
    for caveat in caveats[:6]:
        _rect(slide, Inches(0.38), y + Inches(0.12), Inches(0.055), Inches(0.3), C.AMBER)
        _txt(slide, caveat,
             Inches(0.62), y, Inches(12.3), Inches(0.58),
             size=13.5, colour=C.CHARCOAL)
        y += Inches(0.93)

    _footer(slide, source_note)
    _slide_number(slide, slide_num, total)


def slide_back_cover(prs: Presentation, cfg: DeckConfig) -> None:
    """Minimal back cover — contact / next steps."""
    slide = prs.slides.add_slide(_blank_layout(prs))
    _bg(slide, C.NAVY)
    _rect(slide, Inches(0), Inches(0), Inches(0.22), H, C.TEAL)

    _txt(slide, "Questions or next steps",
         Inches(0.6), Inches(2.0), Inches(10), Inches(0.8),
         size=28, bold=True, colour=C.WHITE)
    _rect(slide, Inches(0.6), Inches(2.95), Inches(6), Inches(0.055), C.TEAL)

    if cfg.contact:
        _txt(slide, cfg.contact,
             Inches(0.6), Inches(3.1), Inches(10), Inches(0.55),
             size=14, colour=C.COVER_SUB)

    _txt(slide, cfg.classification,
         Inches(0.6), Inches(4.2), Inches(10), Inches(0.4),
         size=11, colour=C.SLATE)


# ── Deck assembly ─────────────────────────────────────────────────────────────

@dataclass
class DeckConfig:
    title: str
    subtitle: str
    organisation: str = ""
    prepared_by: str = "Analytics & Insights"
    prepared_date: str = ""
    classification: str = "CONFIDENTIAL — FOR INTERNAL USE ONLY"
    contact: str = ""
    source_note: str = ""

    # Content
    exec_summary_bullets: list[str] = field(default_factory=list)
    chart_slides: list[dict] = field(default_factory=list)
    key_number_slides: list[dict] = field(default_factory=list)
    two_col_slides: list[dict] = field(default_factory=list)
    action_cards: list[dict] = field(default_factory=list)
    caveats: list[str] = field(default_factory=list)
    table_slides: list[dict] = field(default_factory=list)


def build_deck(cfg: DeckConfig, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    if not cfg.prepared_date:
        cfg.prepared_date = datetime.date.today().strftime("%d %B %Y")

    # Count total slides for numbering
    total = 2  # cover + back
    if cfg.exec_summary_bullets:
        total += 1
    total += len(cfg.key_number_slides)
    total += len(cfg.chart_slides)
    total += len(cfg.two_col_slides)
    total += len(cfg.table_slides)
    if cfg.action_cards:
        total += 1
    if cfg.caveats:
        total += 1

    n = 0  # slide counter (cover not numbered)

    slide_cover(prs, cfg)

    if cfg.exec_summary_bullets:
        n += 1
        slide_exec_summary(prs, cfg.exec_summary_bullets, cfg.source_note, n, total - 2)

    for item in cfg.key_number_slides:
        n += 1
        slide_key_number(
            prs,
            title=item["title"],
            subtitle=item.get("subtitle", ""),
            hero_number=item["hero_number"],
            hero_label=item["hero_label"],
            supporting_points=item.get("points", []),
            source_note=item.get("source_note", cfg.source_note),
            slide_num=n, total=total - 2,
        )

    for item in cfg.chart_slides:
        n += 1
        slide_chart(
            prs,
            title=item["title"],
            subtitle=item.get("subtitle", ""),
            chart_path=item["chart_path"],
            annotation_lines=item.get("annotation", []),
            source_note=item.get("source_note", cfg.source_note),
            slide_num=n, total=total - 2,
        )

    for item in cfg.two_col_slides:
        n += 1
        slide_two_col(
            prs,
            title=item["title"],
            subtitle=item.get("subtitle", ""),
            left_heading=item["left_heading"],
            left_items=item["left_items"],
            right_heading=item["right_heading"],
            right_items=item["right_items"],
            source_note=item.get("source_note", cfg.source_note),
            slide_num=n, total=total - 2,
        )

    for item in cfg.table_slides:
        n += 1
        slide_table(
            prs,
            title=item["title"],
            subtitle=item.get("subtitle", ""),
            headers=item["headers"],
            rows=item["rows"],
            col_widths=item["col_widths"],
            source_note=item.get("source_note", cfg.source_note),
            slide_num=n, total=total - 2,
        )

    if cfg.action_cards:
        n += 1
        slide_action_cards(
            prs,
            title="Implications & Recommended Actions",
            subtitle="Priority areas and suggested next steps",
            cards=cfg.action_cards,
            source_note=cfg.source_note,
            slide_num=n, total=total - 2,
        )

    if cfg.caveats:
        n += 1
        slide_caveats(prs, cfg.caveats, cfg.source_note, n, total - 2)

    slide_back_cover(prs, cfg)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Deck saved: {out_path}  ({n + 2} slides)")


# ── CLI: build from JSON config ───────────────────────────────────────────────

def build_from_json(config_path: str) -> None:
    with open(config_path) as f:
        raw = json.load(f)

    cfg = DeckConfig(**{k: v for k, v in raw.items() if k != "output_path"})
    out = Path(raw.get("output_path", "outputs/deck.pptx"))
    build_deck(cfg, out)


# ── Salary sacrifice demo deck (standalone, no JSON needed) ──────────────────

def build_salary_sacrifice_deck() -> None:
    """Rebuild the salary sacrifice analysis deck with the new design system."""
    cfg = DeckConfig(
        title="The Age Gap Is a\nParticipation Problem",
        subtitle="Salary Sacrifice Behaviour Across Age Bands — Superannuation Member Analysis",
        organisation="Analytics & Insights",
        prepared_by="Analytics & Insights",
        prepared_date="29 April 2026",
        contact="Contact your Analytics & Insights team for questions or to run further analysis.",
        source_note="Source: super_members_2026Q1.csv  |  n=5,000 members  |  Snapshot: 31 March 2026  |  Analytics & Insights",

        exec_summary_bullets=[
            "Participation climbs from 9% (age 18–24) to 64% (age 55–64) — the gap is who opts in, not how much they contribute.",
            "Age is a stronger predictor of participation than salary: standardised coefficient 0.110 vs 0.048 for log salary.",
            "Among participants, annual sacrifice amounts are nearly flat across all age bands ($6,024–$7,396 median).",
            "High-earners aged 55–64 participate less than mid-earners — likely a concessional cap effect worth investigating.",
            "Biggest opportunity: 718 non-participating 35–44 year-olds at the inflection point of the participation curve.",
        ],

        key_number_slides=[
            {
                "title": "Participation Is the Lever — Not Contribution Amounts",
                "subtitle": "The ratio of non-participants to participants is the defining gap across all age cohorts",
                "hero_number": "48.7%",
                "hero_label": "Overall fund participation rate\n(2,434 of 5,000 members)",
                "points": [
                    "18–24 band: 9.2% participation (31 of 338 members)",
                    "35–44 band: 48.0% participation — the inflection point",
                    "55–64 band: 63.5% participation — the ceiling",
                    "Conditional median amounts: flat $6,024–$7,396 across all bands",
                    "Strategy implication: focus on activation, not on lifting amounts among existing participants",
                    "718 non-participating 35–44 members represent the single largest near-term opportunity",
                ],
            }
        ],

        chart_slides=[
            {
                "title": "Participation Rate Climbs Steadily With Age",
                "subtitle": "From 9.2% at age 18–24 to 63.5% at 55–64 — a 54-percentage-point gap driven by life stage, not salary alone",
                "chart_path": "outputs/2026-04-29_participation_rate_by_age_band.png",
                "annotation": [
                    "9.2%\n18–24 year-olds participate\n(n=31 of 338)",
                    "63.5%\n55–64 year-olds participate\n(n=475 of 748)",
                    "Age rank coefficient (OLS): 0.110\nLog salary coefficient: 0.048\nAge is 2.3× more predictive.",
                    "Fund average: 48.7% — below the APRA industry benchmark of ~55% for funds of this size.",
                ],
            },
            {
                "title": "Conditional Amounts Are Consistent Once Members Opt In",
                "subtitle": "Among participants, median annual salary sacrifice ranges narrowly — $6,024 to $7,396 regardless of age",
                "chart_path": "outputs/2026-04-29_conditional_amount_by_age_band.png",
                "annotation": [
                    "Median: $6,024–$7,396 across all bands",
                    "18–24 median ($7,396) is based on n=31 — treat as indicative only.",
                    "IQR across all bands: approx. $3,700–$9,100. Wide variation that age does not explain.",
                    "The activation gap is not a generosity gap. Members who start contributing do so at a meaningful rate regardless of age.",
                ],
            },
            {
                "title": "Salary Matters More for Younger Members, Less for Older Ones",
                "subtitle": "Within each salary quartile, older members still participate at higher rates — life-stage effects dominate for 45+",
                "chart_path": "outputs/2026-04-29_participation_by_age_salary_quartile.png",
                "annotation": [
                    "18–34 bands: salary is the gating factor. Q4 earners (top 25%) reach 33–61% participation; Q1 earners near zero.",
                    "45+ bands: high participation even in Q1/Q2 salary — proximity to retirement is the driver.",
                    "Anomaly: 55–64 Q4 earners participate less than Q2/Q3. Likely concessional cap ($30K) becoming binding for high salaries.",
                    "Regression R²=9.5% — large unexplained variation remains. Financial advice, family obligations, and anchoring effects are plausible omitted variables.",
                ],
            },
            {
                "title": "Younger Participants Stretch Further Proportionally",
                "subtitle": "Sacrifice as % of salary declines with age — as salaries rise, the $30K concessional cap constrains effective rates",
                "chart_path": "outputs/2026-04-29_sacrifice_pct_salary_by_age.png",
                "annotation": [
                    "18–24: 7.0% of salary\n25–34: 6.6%\n35–44: 6.3%\n45–54: 6.1%\n55–64: 5.7%\n65+: 5.1%",
                    "Younger participants are contributing more relative to income — a strong signal of motivated behaviour.",
                    "Older members are constrained by the $30K concessional cap (incl. employer SG at 11.5% FY25).",
                    "A 55-year-old on $150K has employer SG of $17,250 — only $12,750 of CC headroom before hitting the cap.",
                ],
            },
        ],

        two_col_slides=[
            {
                "title": "Engagement Signals: What Predicts Participation and What Does Not",
                "subtitle": "App usage and login recency show no positive signal — life-stage factors dominate",
                "left_heading": "Factors associated with higher participation",
                "left_items": [
                    "Age / life stage — strongest predictor across all cohorts",
                    "Salary level — gating factor especially for 18–34 band",
                    "Tenure — participants have slightly higher tenure in every band (β=0.028)",
                    "High Growth investment option — 16.2% vs 7.1% participation in 18–24 (likely selection effect)",
                ],
                "right_heading": "Factors with no predictive signal",
                "right_items": [
                    "App installed — members with app installed participate at lower rates in several bands",
                    "Login recency — median days since login nearly identical for participants vs non-participants",
                    "Digital engagement does not predict salary sacrifice behaviour in this dataset",
                    "Implication: do not assume engaged-app members are more likely to convert",
                ],
            }
        ],

        action_cards=[
            {
                "label": "Age 18–34\n1,125 non-participants",
                "metric": "Current participation: 9–27%",
                "body": "Salary milestone triggers.\n\nTarget members crossing the ~$85K threshold. At lower salaries, focus on habit formation and small starts rather than maximum contribution. Financial education over persuasion.",
            },
            {
                "label": "Age 35–44\n718 non-participants",
                "metric": "Current participation: 48%",
                "body": "Biggest absolute opportunity.\n\nPeer comparison messaging: half their cohort already participates. Default opt-in nudge testing recommended. This cohort has the salary to contribute and the time horizon to benefit most.",
            },
            {
                "label": "Age 55–64\n273 non-participants",
                "metric": "Median salary: $129,200",
                "body": "Investigate before assuming disengagement.\n\nHigh-income non-participants may use alternative strategies (personal deductible contributions). A member service outreach call may outperform digital nudge for this cohort.",
            },
            {
                "label": "Age 65+\n175 active participants",
                "metric": "Near CC cap risk",
                "body": "Concessional cap guidance.\n\nMany are near the $30K annual cap including employer SG. Proactive guidance on carry-forward unused concessional contributions (where TSB < $500K) may unlock additional tax-effective savings.",
            },
        ],

        caveats=[
            "Cross-sectional snapshot (31/03/2026) only — causality cannot be established; prior behaviour not captured.",
            "Salary nulls: 100 members (2%) missing salary data, excluded from salary-stratified analyses. May bias salary results if non-random.",
            "18–24 participant group is small (n=31) — conditional amounts for this band are indicative only; do not over-index.",
            "No adviser relationship data — an omitted variable that likely explains a portion of the unexplained participation variance.",
            "Concessional cap interaction not explicitly modelled — the $30K cap (incl. employer SG) creates a ceiling effect in older, higher-income bands.",
            "Regression model R²=9.5% — significant unexplained variation; findings are descriptive, not causal.",
        ],
    )

    out = Path("outputs/2026-04-29_salary_sacrifice_report_v2.pptx")
    build_deck(cfg, out)


if __name__ == "__main__":
    if len(sys.argv) > 1:
        build_from_json(sys.argv[1])
    else:
        # Default: rebuild the salary sacrifice demo deck
        build_salary_sacrifice_deck()
